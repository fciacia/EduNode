"""Tests for building a downloadable PowerPoint deck from slide data."""
from io import BytesIO

from pptx import Presentation

from core.pptx_export import build_pptx

SLIDES = [
    {"title": "Photosynthesis", "bullets": ["Plants make food", "Uses sunlight"],
     "notes": "Intro to how plants feed.", "emoji": "🌿"},
    {"title": "Summary", "bullets": ["Light becomes food"], "notes": "", "emoji": ""},
]


def _open(blob):
    return Presentation(BytesIO(blob))


def test_build_pptx_one_slide_per_item():
    prs = _open(build_pptx("photosynthesis", SLIDES))
    assert len(prs.slides) == len(SLIDES)


def test_build_pptx_has_titles_and_bullets():
    prs = _open(build_pptx("photosynthesis", SLIDES))
    texts = ["\n".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
             for s in prs.slides]
    assert "Photosynthesis" in texts[0]
    assert "Plants make food" in texts[0]
    assert "Uses sunlight" in texts[0]


def test_build_pptx_puts_notes_in_speaker_notes():
    prs = _open(build_pptx("photosynthesis", SLIDES))
    assert "Intro to how plants feed." in prs.slides[0].notes_slide.notes_text_frame.text


def test_build_pptx_returns_nonempty_bytes():
    blob = build_pptx("x", SLIDES)
    assert isinstance(blob, (bytes, bytearray)) and len(blob) > 1000


def test_build_pptx_embeds_matching_local_image(tmp_path):
    from PIL import Image
    Image.new("RGB", (64, 48), "green").save(tmp_path / "green_leaf.png")
    deck = [{"title": "Leaf", "bullets": ["green"], "image_query": "green leaf", "emoji": "🌿"}]
    prs = _open(build_pptx("leaf", deck, media_dir=tmp_path))
    pics = [sh for sh in prs.slides[0].shapes if sh.shape_type == 13]  # 13 = PICTURE
    assert len(pics) == 1


def test_build_pptx_emoji_fallback_when_no_image(tmp_path):
    deck = [{"title": "Sun", "bullets": ["hot"], "image_query": "nonexistent thing", "emoji": "☀️"}]
    prs = _open(build_pptx("sun", deck, media_dir=tmp_path))
    pics = [sh for sh in prs.slides[0].shapes if sh.shape_type == 13]
    text = "\n".join(sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame)
    assert not pics and "☀️" in text
