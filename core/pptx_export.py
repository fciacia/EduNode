"""
core/pptx_export.py
===================
Render a generated slide deck into a standard PowerPoint (.pptx) file so
students can download it and study offline in PowerPoint / LibreOffice / Google
Slides.

The export mirrors the on-screen deck: per-slide accent colour (bar, title,
bullets), a tinted media tile holding the slide's diagram (or its emoji when no
local image matches), a faint slide-number watermark, and the narration text as
speaker notes.

build_pptx(topic, slides, media_dir=None) -> bytes
"""
from __future__ import annotations

from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from core.media import find_media, load_raster

# Same accent palette the web deck cycles through: (accent, soft tint).
PALETTE = [
    ("6366F1", "EEF0FE"), ("0EA5E9", "E5F6FE"), ("10B981", "E6F8F1"),
    ("F59E0B", "FEF3E0"), ("EC4899", "FDEAF4"), ("8B5CF6", "F1ECFE"),
]

SW, SH = 13.333, 7.5          # 16:9 slide, inches
TEXT_DARK = RGBColor(0x33, 0x33, 0x3A)


def _rgb(hex6: str) -> RGBColor:
    return RGBColor(int(hex6[0:2], 16), int(hex6[2:4], 16), int(hex6[4:6], 16))


def _mix(hex6: str, t: float) -> RGBColor:
    """Blend a hex colour t-of-the-way toward white (0=colour, 1=white)."""
    r, g, b = (int(hex6[i:i + 2], 16) for i in (0, 2, 4))
    return RGBColor(*(int(c + (255 - c) * t) for c in (r, g, b)))


def _no_outline(shape) -> None:
    shape.line.fill.background()
    shape.shadow.inherit = False


def build_pptx(topic: str, slides: list[dict], media_dir=None) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    blank = prs.slide_layouts[6]

    for i, s in enumerate(slides):
        accent, tint = PALETTE[i % len(PALETTE)]
        slide = prs.slides.add_slide(blank)

        # Left accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.22), Inches(SH))
        bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(accent); _no_outline(bar)

        # Faint slide-number watermark, top-right
        wm = slide.shapes.add_textbox(Inches(SW - 2.4), Inches(0.15), Inches(2.1), Inches(1.6))
        wp = wm.text_frame.paragraphs[0]; wp.alignment = PP_ALIGN.RIGHT
        run = wp.add_run(); run.text = f"{i + 1:02d}"
        run.font.size = Pt(72); run.font.bold = True; run.font.color.rgb = _mix(accent, 0.78)

        # Tinted media tile (left)
        tx, ty, tw, th = 0.7, 1.5, 4.3, 4.4
        tile = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(tx), Inches(ty), Inches(tw), Inches(th))
        tile.fill.solid(); tile.fill.fore_color.rgb = _rgb(tint); _no_outline(tile)
        _fill_media(slide, s, tx, ty, tw, th, media_dir)

        # Title (accent) + underline
        title = slide.shapes.add_textbox(Inches(5.4), Inches(1.2), Inches(7.4), Inches(1.4))
        tf = title.text_frame; tf.word_wrap = True
        tp = tf.paragraphs[0]
        emoji = (s.get("emoji") or "").strip()
        tr = tp.add_run(); tr.text = f"{emoji} {(s.get('title') or '').strip()}".strip()
        tr.font.size = Pt(34); tr.font.bold = True; tr.font.color.rgb = _rgb(accent)

        underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.45), Inches(2.62), Inches(6.0), Inches(0.05))
        underline.fill.solid(); underline.fill.fore_color.rgb = _mix(accent, 0.45); _no_outline(underline)

        # Bullets with accent dots
        box = slide.shapes.add_textbox(Inches(5.4), Inches(2.9), Inches(7.4), Inches(3.9))
        bt = box.text_frame; bt.word_wrap = True
        bullets = [b for b in (s.get("bullets") or []) if str(b).strip()]
        for j, bullet in enumerate(bullets):
            para = bt.paragraphs[0] if j == 0 else bt.add_paragraph()
            para.space_after = Pt(10)
            dot = para.add_run(); dot.text = "●  "
            dot.font.size = Pt(20); dot.font.color.rgb = _rgb(accent)
            txt = para.add_run(); txt.text = str(bullet).strip()
            txt.font.size = Pt(20); txt.font.color.rgb = TEXT_DARK

        # Narration -> speaker notes
        notes = (s.get("notes") or "").strip()
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    out = BytesIO()
    prs.save(out)
    return out.getvalue()


def _fill_media(slide, s, tx, ty, tw, th, media_dir) -> None:
    """Embed a matching local image inside the tile, else show the emoji."""
    query = (s.get("image_query") or "").strip()
    path = find_media(query, root=media_dir) if query else None
    raster = load_raster(path) if path else None

    if raster is not None:
        try:
            from PIL import Image
            iw, ih = Image.open(raster).size
            raster.seek(0)
            pad = 0.35
            box_w, box_h = tw - 2 * pad, th - 2 * pad
            ar = iw / ih if ih else 1
            if box_w / box_h > ar:          # tile wider than image -> bound by height
                h, w = box_h, box_h * ar
            else:
                w, h = box_w, box_w / ar
            slide.shapes.add_picture(
                raster,
                Inches(tx + (tw - w) / 2), Inches(ty + (th - h) / 2),
                Inches(w), Inches(h),
            )
            return
        except Exception:
            pass

    # Emoji fallback, centred in the tile
    emoji = (s.get("emoji") or "📘").strip() or "📘"
    tb = slide.shapes.add_textbox(Inches(tx), Inches(ty), Inches(tw), Inches(th))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = emoji; run.font.size = Pt(96)
