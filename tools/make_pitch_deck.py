"""
tools/make_pitch_deck.py
========================
Generate the 15-slide ASEAN AI Hackathon pitch deck for Edge as a designed .pptx
— charts, an architecture diagram, stat callouts and a real product screenshot.
All metrics are the real measured values from the evaluation harness.

Run: python -m tools.make_pitch_deck
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# ── palette ──────────────────────────────────────────────────────────────────
PURPLE   = RGBColor(0x5B, 0x3F, 0xE0)
PURPLE_D = RGBColor(0x3B, 0x27, 0xA0)
PANEL    = RGBColor(0xF1, 0xEE, 0xFC)
PANEL2   = RGBColor(0xEA, 0xF6, 0xEE)
DARK     = RGBColor(0x1E, 0x1E, 0x28)
GREY     = RGBColor(0x5A, 0x5A, 0x66)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREEN    = RGBColor(0x1E, 0x8E, 0x3E)
AMBER    = RGBColor(0xB0, 0x60, 0x00)
TEAL     = RGBColor(0x0E, 0x80, 0x9B)

W, H = Inches(13.333), Inches(7.5)
ASSETS = Path("docs/assets")


# ── primitives ───────────────────────────────────────────────────────────────
def _bg(s, color):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = color


def _txt(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    return tf


def _line(p, text, size, color=DARK, bold=False, align=PP_ALIGN.LEFT, italic=False,
          space_after=6):
    p.text = text; p.alignment = align; p.space_after = Pt(space_after)
    r = p.runs[0]
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = "Calibri"
    return p


def _rect(s, l, t, w, h, fill, rounded=True, line_color=None):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                             l, t, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_color:
        shp.line.color.rgb = line_color; shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _arrow(s, l, t, w, h, color=PURPLE):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background()
    a.shadow.inherit = False
    return a


def _header(s, kicker, title):
    _rect(s, Inches(0), Inches(0), Inches(0.18), H, PURPLE, rounded=False)
    tf = _txt(s, Inches(0.6), Inches(0.4), Inches(12.4), Inches(1.1))
    _line(tf.paragraphs[0], kicker.upper(), 13, PURPLE, bold=True, space_after=2)
    _line(tf.add_paragraph(), title, 30, DARK, bold=True)
    # footer
    ft = _txt(s, Inches(0.6), Inches(7.02), Inches(12.1), Inches(0.4))
    _line(ft.paragraphs[0], "Edge — AI Education at the Last Mile   ·   UM-Gunners",
          10, GREY)


def _stat_card(s, l, t, w, big, label, accent=PURPLE, sub=None):
    card = _rect(s, l, t, w, Inches(1.9), PANEL)
    tf = card.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18); tf.margin_top = Inches(0.18)
    _line(tf.paragraphs[0], big, 40, accent, bold=True, space_after=2)
    _line(tf.add_paragraph(), label, 14, DARK, bold=True, space_after=2)
    if sub:
        _line(tf.add_paragraph(), sub, 11, GREY)


def _bullets(s, l, t, w, h, items, size=18):
    tf = _txt(s, l, t, w, h)
    first = True
    for it in items:
        lvl = 0
        if isinstance(it, tuple):
            it, lvl = it
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        _line(p, ("•  " if lvl == 0 else "–  ") + it, size if lvl == 0 else size-3,
              DARK if lvl == 0 else GREY, space_after=9 if lvl == 0 else 4)
        p.level = lvl
    return tf


def _chart(s, l, t, w, h, cats, series, title=None, colors=None):
    cd = CategoryChartData(); cd.categories = cats
    for name, vals in series:
        cd.add_series(name, vals)
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, l, t, w, h, cd)
    ch = gf.chart
    ch.has_legend = len(series) > 1
    if ch.has_legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM; ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(11)
    if title:
        ch.has_title = True; ch.chart_title.text_frame.text = title
        ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
    pal = colors or [PURPLE, GREEN, TEAL]
    for i, plot_series in enumerate(ch.series):
        plot_series.format.fill.solid()
        plot_series.format.fill.fore_color.rgb = pal[i % len(pal)]
    ch.value_axis.has_major_gridlines = False
    try:
        ch.category_axis.tick_labels.font.size = Pt(10)
        ch.value_axis.tick_labels.font.size = Pt(9)
    except Exception:
        pass
    return ch


# ── slides ───────────────────────────────────────────────────────────────────
def s_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, PURPLE)
    _rect(s, Inches(0), Inches(5.9), W, Inches(1.6), PURPLE_D, rounded=False)
    tf = _txt(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(3.4))
    _line(tf.paragraphs[0], "Edge", 80, WHITE, bold=True, space_after=2)
    _line(tf.add_paragraph(), "AI Education at the Last Mile", 32, RGBColor(0xE7,0xE0,0xFF))
    bf = _txt(s, Inches(0.95), Inches(6.05), Inches(11.5), Inches(1.3))
    _line(bf.paragraphs[0], "UM – Gunners   ·   University of Malaya   ·   Malaysia",
          18, WHITE, bold=True, space_after=3)
    _line(bf.add_paragraph(), "Track: AI for Education   ·   ASEAN AI Hackathon 2026",
          14, RGBColor(0xD8,0xCE,0xFF))


def s_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _header(s, "The Challenge", "60M+ rural learners, left offline")
    _stat_card(s, Inches(0.7), Inches(1.7), Inches(3.8), "60M+",
               "rural ASEAN learners excluded", PURPLE, "no internet · no grid · no mother-tongue AI")
    _stat_card(s, Inches(4.75), Inches(1.7), Inches(3.8), "18%",
               "rural internet (Timor-Leste)", TEAL, "25% Myanmar · 28% Laos (ITU 2023)")
    _stat_card(s, Inches(8.8), Inches(1.7), Inches(3.8), "55%",
               "Filipino 15-yos below PISA maths", AMBER, "30% Cambodian G5 read (UNICEF 2019)")
    _bullets(s, Inches(0.7), Inches(4.0), Inches(12.0), Inches(2.6), [
        "Cloud AI (ChatGPT, etc.) is built for the connected world and fails entirely off-grid",
        "The communities with the least connectivity have the least access to quality tutoring",
        "Mother-tongue and minority-language learners are unserved by dominant-language models",
    ])


def s_solution(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _header(s, "The Solution", "An offline AI tutor on a $220 solar hub")
    _bullets(s, Inches(0.7), Inches(1.75), Inches(12.2), Inches(2.0), [
        "A fully offline, solar-powered AI tutoring hub on a Raspberry Pi — broadcasts local Wi-Fi",
        "Students connect from their own phones, no internet; an Agentic RAG pipeline answers in their mother tongue",
        "Curriculum-grounded answers, adaptive quizzes, voice interaction, PDF microcredentials",
    ], size=17)
    feats = [("100% Offline", "no internet, ever", PURPLE),
             ("Mother-tongue", "incl. Iban, Cebuano", TEAL),
             ("Grounded", "0% hallucination", GREEN),
             ("$0.60 / student / yr", "solar, no cloud fees", AMBER)]
    x = Inches(0.7)
    for big, sub, col in feats:
        c = _rect(s, x, Inches(4.35), Inches(2.95), Inches(1.7), PANEL)
        tf = c.text_frame; tf.word_wrap = True; tf.margin_top = Inches(0.2); tf.margin_left = Inches(0.16)
        _line(tf.paragraphs[0], big, 18, col, bold=True, space_after=3)
        _line(tf.add_paragraph(), sub, 12, GREY)
        x += Inches(3.08)


def s_advantage(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _header(s, "Competitive Advantage", "Why Edge, not the cloud")
    # two columns: Cloud AI vs Edge
    _rect(s, Inches(0.7), Inches(1.8), Inches(5.85), Inches(4.7), PANEL)
    _rect(s, Inches(6.8), Inches(1.8), Inches(5.85), Inches(4.7), PANEL2)
    a = _txt(s, Inches(0.95), Inches(2.0), Inches(5.4), Inches(4.3))
    _line(a.paragraphs[0], "Cloud AI (ChatGPT…)", 18, GREY, bold=True, space_after=10)
    for t in ["Needs constant internet", "Recurring subscription cost",
              "Dominant languages only", "Ungrounded — hallucinates",
              "No teacher analytics / governance"]:
        _line(a.add_paragraph(), "✕  " + t, 14, GREY, space_after=8)
    b = _txt(s, Inches(7.05), Inches(2.0), Inches(5.4), Inches(4.3))
    _line(b.paragraphs[0], "Edge", 18, GREEN, bold=True, space_after=10)
    for t in ["Runs 100% offline on solar", "≈$0.60/student/yr, no fees",
              "Mother-tongue incl. low-resource", "Curriculum-grounded, 0% false-grounding",
              "Teacher dashboard + full governance"]:
        _line(b.add_paragraph(), "✓  " + t, 14, DARK, space_after=8)


def s_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _header(s, "Technical Architecture", "A 4-agent offline Agentic RAG pipeline")
    # input
    _rect(s, Inches(0.7), Inches(2.6), Inches(1.9), Inches(1.6), TEAL)
    t = _txt(s, Inches(0.75), Inches(2.7), Inches(1.8), Inches(1.4), MSO_ANCHOR.MIDDLE)
    _line(t.paragraphs[0], "Voice / text\nquery\n(dialect)", 12, WHITE, bold=True, align=PP_ALIGN.CENTER)
    # 4 agent boxes
    agents = ["Translation\nAgent", "Personalised\nContext", "Pedagogical\nReasoning", "Verification\nAgent"]
    x = Inches(3.0)
    for i, name in enumerate(agents):
        _arrow(s, x - Inches(0.42), Inches(3.15), Inches(0.4), Inches(0.5), PURPLE)
        box = _rect(s, x, Inches(2.6), Inches(1.75), Inches(1.6), PURPLE if i % 2 == 0 else PURPLE_D)
        tf = _txt(s, x, Inches(2.7), Inches(1.75), Inches(1.4), MSO_ANCHOR.MIDDLE)
        _line(tf.paragraphs[0], f"{i+1}. {name}", 12, WHITE, bold=True, align=PP_ALIGN.CENTER)
        x += Inches(2.18)
    _arrow(s, x - Inches(0.42), Inches(3.15), Inches(0.4), Inches(0.5), GREEN)
    _rect(s, x, Inches(2.6), Inches(1.95), Inches(1.6), GREEN)
    o = _txt(s, x, Inches(2.7), Inches(1.95), Inches(1.4), MSO_ANCHOR.MIDDLE)
    _line(o.paragraphs[0], "Grounded\nmother-tongue\nanswer + citations", 11, WHITE, bold=True, align=PP_ALIGN.CENTER)
    # data layer + hardware
    dl = _txt(s, Inches(0.7), Inches(4.55), Inches(12.0), Inches(0.9))
    _line(dl.paragraphs[0], "Knowledge: ChromaDB (cross-lingual vectors over OER PDFs) · SQLite (profiles, progress, audit)",
          13, GREY, align=PP_ALIGN.CENTER)
    _rect(s, Inches(0.7), Inches(5.5), Inches(12.0), Inches(1.0), PANEL)
    hw = _txt(s, Inches(0.8), Inches(5.6), Inches(11.8), Inches(0.8), MSO_ANCHOR.MIDDLE)
    _line(hw.paragraphs[0], "Hardware (100% offline):  Solar panel → battery → Raspberry Pi → Wi-Fi router → students' phones    ·    two-node split for memory-constrained Pis",
          13, DARK, bold=True, align=PP_ALIGN.CENTER)


def s_models(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _header(s, "AI Approach", "Open-weight models, chosen for the edge")
    cards = [
        ("Phi-3 Mini", "4-bit GGUF · Ollama", "Reasoning & tutoring — small enough for on-device inference", PURPLE),
        ("MiniLM-L12 (multilingual)", "paraphrase-multilingual", "Cross-lingual retrieval — the key to the bias fix", TEAL),
        ("NLLB-200-distilled-600M", "+ glossary bridge", "NMT across ASEAN languages; bridge for non-NLLB dialects", GREEN),
        ("Whisper.cpp + Piper", "ggml-base", "Offline speech-to-text and text-to-speech", AMBER),
    ]
    coords = [(Inches(0.7), Inches(1.8)), (Inches(6.8), Inches(1.8)),
              (Inches(0.7), Inches(4.2)), (Inches(6.8), Inches(4.2))]
    for (l, t), (name, tag, desc, col) in zip(coords, cards):
        _rect(s, l, t, Inches(5.85), Inches(2.15), PANEL)
        tf = _txt(s, l + Inches(0.25), t + Inches(0.18), Inches(5.4), Inches(1.8))
        _line(tf.paragraphs[0], name, 18, col, bold=True, space_after=2)
        _line(tf.add_paragraph(), tag, 12, GREY, italic=True, space_after=6)
        _line(tf.add_paragraph(), desc, 13, DARK)
    cap = _txt(s, Inches(0.7), Inches(6.5), Inches(12.0), Inches(0.4))
    _line(cap.paragraphs[0], "Why: open-weight, permissively licensed, ARM-optimised, fully offline.",
          12, GREY, italic=True)


def s_data(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _header(s, "Data Strategy", "Curated, licensed, cleaned")
    _bullets(s, Inches(0.7), Inches(1.8), Inches(12.0), Inches(3.2), [
        "Curriculum: OpenStax (CC BY 4.0) + national OER from KPM, DepEd, Kemendikbud + team-authored lessons",
        "Translation references: open VynerCK/Iban-language-data + teacher-verified corrections (the Dialect Flywheel)",
        "Cleaning: manual curation against ministry standards · page-level chunking · idempotent ingestion",
        "Integrity: versioned, provenance-tracked content manifest with enforced tamper quarantine",
    ])
    lic = _rect(s, Inches(0.7), Inches(5.1), Inches(12.0), Inches(1.3), PANEL)
    tf = lic.text_frame; tf.margin_top = Inches(0.18); tf.margin_left = Inches(0.2)
    _line(tf.paragraphs[0], "Licensing", 14, PURPLE, bold=True, space_after=4)
    _line(tf.add_paragraph(),
          "OpenStax CC BY 4.0   ·   Phi-3 MIT   ·   Piper Apache 2.0   ·   NLLB-200 CC BY-NC 4.0   ·   Iban/Cebuano CC BY 4.0",
          13, DARK)


def s_ethics(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _header(s, "AI Ethics & Responsibility", "Bias: measured, then mitigated")
    _chart(s, Inches(0.6), Inches(1.8), Inches(6.6), Inches(4.7),
           ["English", "B. Melayu", "Thai", "Viet", "Iban"],
           [("Before fix", (100, 75, 50, 50, 88)), ("After cross-lingual", (100, 75, 62, 62, 100))],
           title="Grounding rate by language (same question)")
    _bullets(s, Inches(7.4), Inches(1.9), Inches(5.4), Inches(3.4), [
        "Same question, varied language → English 100% vs Thai 50% grounding = measured bias",
        "Root cause: query translation degraded retrieval",
        "Fix: cross-lingual retrieval → Thai/Viet 62%, Iban fully closed; 0% non-response",
    ], size=14)
    pr = _rect(s, Inches(7.4), Inches(5.4), Inches(5.4), Inches(1.1), PANEL2)
    tf = pr.text_frame; tf.margin_top = Inches(0.14); tf.margin_left = Inches(0.18)
    _line(tf.paragraphs[0], "Privacy & governance", 13, GREEN, bold=True, space_after=3)
    _line(tf.add_paragraph(),
          "RBAC + per-teacher audit · encryption (LUKS + USB) · parental consent · retention & erasure",
          12, DARK)


def s_demo(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _header(s, "Prototype Demonstration", "The AI, grounded and transparent")
    shot = ASSETS / "demo_chat2.png"
    if shot.exists():
        s.shapes.add_picture(str(shot), Inches(0.7), Inches(1.75), height=Inches(5.0))
    _bullets(s, Inches(6.0), Inches(2.0), Inches(6.8), Inches(3.6), [
        "Real grounded answer with page-level citations (3 sources shown)",
        "Green “Verified” confidence badge — transparent about certainty",
        "Auto-generated fraction diagram — visual, mother-tongue ready",
        "Tiered logic: grounded → supplementary → controlled non-response",
        "Mobile-first PWA: students use their own phones, fully offline",
    ], size=15)
    v = _txt(s, Inches(6.0), Inches(5.9), Inches(6.8), Inches(0.6))
    _line(v.paragraphs[0], "▶  Demo video: https://youtu.be/Ms_y6FUbKwU", 14, PURPLE, bold=True)


def s_hurdles(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _header(s, "Technical Hurdles", "Honest reflection — problem → fix")
    pairs = [
        ("Hallucination in low-context queries", "Verification Agent + tiered distance gates"),
        ("Automated LLM judges unreliable (3.4–97%)", "Validate answer correctness with humans"),
        ("Translation “looked weak” (chrF 28)", "Measurement bug — real chrF is 56"),
        ("Query translation degraded retrieval", "Cross-lingual retrieval (closed the bias gap)"),
        ("Memory pressure on a 4 GB Pi", "Admission queue + TTS offload + two-node split"),
    ]
    y = Inches(1.85)
    for prob, fix in pairs:
        _rect(s, Inches(0.7), y, Inches(5.8), Inches(0.85), PANEL)
        _arrow(s, Inches(6.55), y + Inches(0.2), Inches(0.5), Inches(0.45), PURPLE)
        _rect(s, Inches(7.15), y, Inches(5.5), Inches(0.85), PANEL2)
        pf = _txt(s, Inches(0.9), y + Inches(0.04), Inches(5.5), Inches(0.78), MSO_ANCHOR.MIDDLE)
        _line(pf.paragraphs[0], prob, 13, DARK, bold=True)
        ff = _txt(s, Inches(7.35), y + Inches(0.04), Inches(5.2), Inches(0.78), MSO_ANCHOR.MIDDLE)
        _line(ff.paragraphs[0], fix, 13, GREEN, bold=True)
        y += Inches(0.98)


def s_metrics(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _header(s, "Accuracy & Efficiency Metrics", "Measured, not claimed")
    _stat_card(s, Inches(0.7), Inches(1.7), Inches(2.85), "100%", "Retrieval recall", PURPLE, "right source always found")
    _stat_card(s, Inches(3.75), Inches(1.7), Inches(2.85), "0%", "Hallucination", GREEN, "no false grounding")
    _stat_card(s, Inches(6.8), Inches(1.7), Inches(2.85), "~97%", "Concept coverage", TEAL, "answers contain key facts")
    _stat_card(s, Inches(9.85), Inches(1.7), Inches(2.85), "56", "Iban chrF", AMBER, "decent; flywheel grows it")
    _chart(s, Inches(0.7), Inches(3.85), Inches(6.4), Inches(2.7),
           ["Recall", "Precision", "Concept cov.", "Non-resp."],
           [("score %", (100, 79, 97, 0))], title="RAG quality (%)", colors=[PURPLE])
    note = _txt(s, Inches(7.4), Inches(3.95), Inches(5.4), Inches(2.6))
    _line(note.paragraphs[0], "Honest caveat", 14, AMBER, bold=True, space_after=5)
    _line(note.add_paragraph(),
          "We report concept coverage, not a single automated “accuracy %”: small LLM judges scored "
          "the same answers 3.4–97%, so answer correctness is being validated by human raters.",
          13, DARK, space_after=8)
    _line(note.add_paragraph(),
          "Load test (dev node): 0 errors under the admission queue; throughput plateaus at the "
          "single-stream SLM rate. On-Pi benchmarking is the next step.", 13, GREY)


def s_scalability(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _header(s, "Scalability Roadmap", "From one school to the region")
    steps = [("1 Hub", "~50 connected students,\nadmission-controlled"),
             ("1 School", "two-node split for\nhigher concurrency"),
             ("Many Schools", "decentralised content\nmanifest + USB sync"),
             ("ASEAN-wide", "per-hub language\nconfiguration")]
    x = Inches(0.7)
    for i, (big, sub) in enumerate(steps):
        if i: _arrow(s, x - Inches(0.5), Inches(2.9), Inches(0.45), Inches(0.5), PURPLE)
        c = _rect(s, x, Inches(2.2), Inches(2.7), Inches(1.9), PANEL)
        tf = _txt(s, x, Inches(2.4), Inches(2.7), Inches(1.6), MSO_ANCHOR.MIDDLE)
        _line(tf.paragraphs[0], big, 20, PURPLE, bold=True, align=PP_ALIGN.CENTER, space_after=4)
        _line(tf.add_paragraph(), sub, 12, DARK, align=PP_ALIGN.CENTER)
        x += Inches(3.15)
    _bullets(s, Inches(0.7), Inches(4.6), Inches(12.0), Inches(1.8), [
        "No internet needed to update or scale — encrypted USB sneakernet deploys content & DB updates",
        "Ministries, NGOs, schools and teachers publish curriculum independently (provenance-tracked)",
    ], size=15)


def s_impact(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _header(s, "Impact Assessment", "Aligned with the UN SDGs")
    sdgs = [("SDG 4", "Quality Education", "curriculum-grounded tutoring at the last mile", PURPLE),
            ("SDG 10", "Reduced Inequalities", "mother-tongue access for underserved learners", TEAL),
            ("SDG 9", "Infrastructure & Innovation", "resilient off-grid education infrastructure", GREEN)]
    x = Inches(0.7)
    for code, name, desc, col in sdgs:
        c = _rect(s, x, Inches(2.0), Inches(3.85), Inches(2.6), PANEL)
        tf = _txt(s, x + Inches(0.25), Inches(2.25), Inches(3.4), Inches(2.2))
        _line(tf.paragraphs[0], code, 30, col, bold=True, space_after=2)
        _line(tf.add_paragraph(), name, 16, DARK, bold=True, space_after=6)
        _line(tf.add_paragraph(), desc, 13, GREY)
        x += Inches(4.05)
    b = _txt(s, Inches(0.7), Inches(5.0), Inches(12.0), Inches(1.4))
    _line(b.paragraphs[0],
          "Regional resilience: learning continuity off-grid and in crises; preservation of minority "
          "languages via the ASEAN Dialect Flywheel — a public good of minority-language AI infrastructure.",
          15, DARK)


def s_future(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, WHITE)
    _header(s, "Future Roadmap", "Next steps after the hackathon")
    _bullets(s, Inches(0.7), Inches(1.9), Inches(12.0), Inches(4.6), [
        "Micro-pilot in rural schools → measured learning outcomes (pre/post, small-n)",
        "On-Pi benchmarking + two-node deployment validation",
        "Fluent-speaker translation audit → sentence-level chrF; grow the dialect-glossary flywheel",
        "Stronger embedder / fine-tuned MT as teacher-verified data accrues",
        "Hardening: encryption rollout, parental-consent workflow, retention automation",
    ], size=18)


def s_team(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, PURPLE)
    tf = _txt(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(1.2))
    _line(tf.paragraphs[0], "Team & Contact", 34, WHITE, bold=True, space_after=2)
    _line(tf.add_paragraph(), "UM – Gunners · University of Malaya · Malaysia", 18, RGBColor(0xE0,0xD8,0xFF))
    x = Inches(0.9)
    for n in range(3):
        c = _rect(s, x, Inches(2.4), Inches(3.7), Inches(2.6), WHITE)
        cf = _txt(s, x + Inches(0.25), Inches(2.65), Inches(3.2), Inches(2.1))
        _line(cf.paragraphs[0], f"[ Member {n+1} ]", 18, PURPLE, bold=True, space_after=4)
        _line(cf.add_paragraph(), "Role / contribution", 13, GREY)
        x += Inches(3.95)
    lf = _txt(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(1.0))
    _line(lf.paragraphs[0], "Repository:  github.com/fciacia/EduNode", 16, WHITE, bold=True, space_after=3)
    _line(lf.add_paragraph(), "Demo:  youtu.be/Ms_y6FUbKwU   ·   [ add university logo ]", 14, RGBColor(0xE0,0xD8,0xFF))


def build() -> Path:
    prs = Presentation(); prs.slide_width, prs.slide_height = W, H
    for fn in (s_title, s_problem, s_solution, s_advantage, s_architecture, s_models,
               s_data, s_ethics, s_demo, s_hurdles, s_metrics, s_scalability, s_impact,
               s_future, s_team):
        fn(prs)
    out = Path("docs/Edge_PitchDeck.pptx"); out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    return out


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path} ({len(Presentation(path).slides._sldIdLst)} slides)")
